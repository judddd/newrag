"""Document processing module with LangChain integration"""

import hashlib
import json
import os
import tempfile
import zipfile
import concurrent.futures
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional

import structlog
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredHTMLLoader,
    CSVLoader,
    UnstructuredExcelLoader,
)
from langchain_core.documents import Document

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

from src.config import config
from src.models import VisionModel
from src.vlm_extractor import VLMPageExtractor

logger = structlog.get_logger(__name__)


def detect_garbled_text(text: str, threshold: float = 0.15) -> tuple[bool, float]:
    """
    检测文本是否包含过多乱码
    
    Args:
        text: 要检测的文本
        threshold: 乱码率阈值（默认 15%）
    
    Returns:
        (is_garbled, garbled_ratio): 是否乱码，乱码率
    """
    if not text or len(text) == 0:
        return False, 0.0
    
    # 统计各类字符
    total_chars = len(text)
    garbled_count = 0
    
    for char in text:
        code_point = ord(char)
        
        # 检测常见乱码字符范围
        # 1. 控制字符（除了常见的空白符）
        if 0x0000 <= code_point <= 0x001F and char not in ['\n', '\r', '\t']:
            garbled_count += 1
        # 2. 私用区
        elif 0xE000 <= code_point <= 0xF8FF:
            garbled_count += 1
        # 3. 特殊符号区的异常字符
        elif code_point in [0xFFFD, 0xFFFE, 0xFFFF]:  # 替换字符
            garbled_count += 1
        # 4. 连续的框框字符 (tofu/豆腐块)
        elif 0x2580 <= code_point <= 0x259F:
            garbled_count += 1
    
    garbled_ratio = garbled_count / total_chars if total_chars > 0 else 0.0
    is_garbled = garbled_ratio > threshold
    
    return is_garbled, garbled_ratio


class DocumentProcessor:
    """Document processor with multi-format support and metadata extraction"""

    def __init__(self, config_override: Optional[Dict[str, Any]] = None):
        """
        Initialize document processor
        
        Args:
            config_override: Optional config dict to override global config
        """
        # 从 config.yaml 读取所有配置
        self.config = config_override or config.processing_config
        self.text_splitting_config = config.text_splitting_config
        self.metadata_config = config.metadata_config
        
        # Initialize text splitter
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.text_splitting_config.get('chunk_size', 500),
            chunk_overlap=self.text_splitting_config.get('chunk_overlap', 50),
            separators=self.text_splitting_config.get('separators', ["\n\n", "\n", " ", ""]),
            length_function=len,
        )
        
        # Initialize vision model if enabled
        self.vision_model = None
        if self.config.get('extract_images', False):
            try:
                self.vision_model = VisionModel()
            except Exception as e:
                logger.warning("vision_model_init_failed", error=str(e))
        
        # Initialize VLM page extractor
        self.vlm_extractor = None
        try:
            self.vlm_extractor = VLMPageExtractor()
            logger.info("vlm_extractor_initialized")
        except Exception as e:
            logger.warning("vlm_extractor_init_failed", error=str(e))
        
        logger.info("document_processor_initialized")
    
    def load_document(self, file_path: str, processed_json_dir: Optional[str] = None) -> List[Document]:
        """
        Load document using appropriate loader
        
        Args:
            file_path: Path to document file
            processed_json_dir: Optional path to directory containing complete_document.json
        
        Returns:
            List of LangChain Document objects
        """
        file_path = Path(file_path)
        file_ext = file_path.suffix.lower()
        
        try:
            if file_ext == '.pdf':
                # Check if we have pre-processed JSON (from adaptive OCR pipeline)
                if processed_json_dir:
                    complete_json = Path(processed_json_dir) / "complete_document.json"
                    if complete_json.exists():
                        logger.info("loading_from_preprocessed_json", json_path=str(complete_json))
                        return self._load_from_complete_json(complete_json)
                    else:
                        logger.debug("complete_json_not_found", expected_path=str(complete_json))
                loader = PyPDFLoader(str(file_path))
                documents = loader.load()
                
                logger.info(
                    "pdf_loaded",
                    file_path=str(file_path),
                    num_pages=len(documents)
                )
                
                # Process pages concurrently
                processed_documents = [None] * len(documents)
                
                def process_single_page(idx, doc):
                    try:
                        # Detect page content type
                        page_type = self.detect_page_content_type(doc.page_content)
                        
                        # Check if page has garbled text or is a drawing
                        is_garbled, garbled_ratio = detect_garbled_text(doc.page_content) if doc.page_content else (False, 0.0)
                        content_len = len(doc.page_content.strip()) if doc.page_content else 0
                        
                        # Use VLM if:
                        # 1. Garbled text
                        # 2. Very short content (< 50 chars)
                        # 3. Drawing (ALWAYS use vision for drawings to capture spatial layout)
                        # 4. Table AND content is not rich (tables with rich text are usually extractable)
                        # 5. Explicitly configured to convert tables
                        is_rich_text = content_len > 800
                        
                        needs_vlm = (
                            is_garbled or 
                            content_len < 50 or
                            page_type == 'drawing' or  # Always use vision for drawings
                            (page_type == 'table' and not is_rich_text and self.config.get('convert_to_image_for_tables', True))
                        )
                        
                        logger.info(
                            "pdf_page_analysis",
                            page=idx + 1,
                            page_type=page_type,
                            is_garbled=is_garbled,
                            garbled_ratio=f"{garbled_ratio:.2%}",
                            needs_vlm=needs_vlm,
                            content_length=content_len
                        )
                        
                        # Use VLM extraction if needed
                        if needs_vlm and self.vlm_extractor and PDF2IMAGE_AVAILABLE:
                            try:
                                logger.info("using_vlm_for_page", page=idx + 1)
                                
                                # Convert page to image
                                image_path = self._convert_pdf_page_to_image(file_path, idx)
                                
                                try:
                                    # Extract with VLM
                                    page_json = self.vlm_extractor.extract_page_content(image_path, page_type)
                                    
                                    # Convert JSON to searchable text
                                    doc.page_content = self._flatten_to_searchable_text(page_json)
                                    doc.metadata['page_json'] = page_json
                                    doc.metadata['page_type'] = page_type
                                    doc.metadata['extraction_method'] = 'vlm'
                                    
                                    logger.info(
                                        "vlm_extraction_success",
                                        page=idx + 1,
                                        num_equipment=len(page_json.get('equipment', [])),
                                        num_components=len(page_json.get('components', []))
                                    )
                                finally:
                                    # Clean up temp image
                                    try:
                                        os.unlink(image_path)
                                    except:
                                        pass
                            
                            except Exception as e:
                                logger.error(
                                    "vlm_extraction_failed",
                                    page=idx + 1,
                                    error=str(e)
                                )
                                # Keep original content
                                doc.metadata['page_type'] = page_type
                                doc.metadata['extraction_method'] = 'text'
                        else:
                            # Use text content
                            doc.metadata['page_type'] = page_type
                            doc.metadata['extraction_method'] = 'text'
                        
                        return idx, doc
                    except Exception as e:
                        logger.error("page_processing_error", page=idx+1, error=str(e))
                        return idx, doc

                # Use ThreadPoolExecutor for concurrent processing
                # Limit workers to avoid overloading LM Studio or Memory
                max_workers = self.config.get('max_workers', 4)
                with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
                    futures = [executor.submit(process_single_page, i, doc) for i, doc in enumerate(documents)]
                    for future in concurrent.futures.as_completed(futures):
                        idx, processed_doc = future.result()
                        processed_documents[idx] = processed_doc
                
                documents = processed_documents
                
                logger.info(
                    "document_loaded",
                    file_path=str(file_path),
                    file_ext=file_ext,
                    num_documents=len(documents)
                )
                
                for idx, doc in enumerate(documents):
                    content_length = len(doc.page_content) if doc.page_content else 0
                    is_garbled, garbled_ratio = detect_garbled_text(doc.page_content) if doc.page_content else (False, 0.0)
                    
                    logger.info(
                        "document_page_info",
                        page_index=idx,
                        content_length=content_length,
                        has_content=bool(doc.page_content and doc.page_content.strip()),
                        is_garbled=is_garbled,
                        garbled_ratio=f"{garbled_ratio:.2%}",
                        content_preview=doc.page_content[:200] if doc.page_content else "EMPTY"
                    )
                
                return documents
            elif file_ext in ['.docx', '.doc']:
                # Check if we have pre-processed JSON (e.g. from new OCR pipeline)
                if processed_json_dir:
                    complete_json = Path(processed_json_dir) / "complete_document.json"
                    if complete_json.exists():
                        logger.info("loading_docx_from_preprocessed_json", json_path=str(complete_json))
                        return self._load_from_complete_json(complete_json)
                    else:
                        # Try adaptive OCR json as fallback
                        complete_adaptive_json = Path(processed_json_dir) / "complete_adaptive_ocr.json"
                        if complete_adaptive_json.exists():
                            logger.info("loading_docx_from_adaptive_json", json_path=str(complete_adaptive_json))
                            return self._load_from_complete_json(complete_adaptive_json)
                
                # Fallback: Try to split Word by pages using page breaks
                try:
                    return self._split_word_by_pages(file_path)
                except Exception as e:
                    logger.warning("word_page_split_failed", error=str(e))
                    # Fall back to standard loading
                    loader = UnstructuredWordDocumentLoader(str(file_path))
            elif file_ext in ['.txt', '.md']:
                loader = TextLoader(str(file_path), encoding='utf-8')
            elif file_ext in ['.html', '.htm']:
                loader = UnstructuredHTMLLoader(str(file_path))
            elif file_ext == '.csv':
                loader = CSVLoader(str(file_path))
            elif file_ext in ['.xlsx', '.xls']:
                # Process Excel by sheets
                try:
                    return self._process_excel_sheets(file_path)
                except Exception as e:
                    logger.warning("excel_sheet_processing_failed", error=str(e))
                    # Fall back to standard loading
                    loader = UnstructuredExcelLoader(str(file_path))
            elif file_ext == '.pptx':
                # Check if we have pre-processed JSON from PPTX pipeline (same structure as PDF)
                if processed_json_dir:
                    complete_json = Path(processed_json_dir) / "complete_adaptive_ocr.json"
                    if complete_json.exists():
                        logger.info("loading_pptx_from_preprocessed_json", json_path=str(complete_json))
                        # Load PPTX data from complete_adaptive_ocr.json
                        with open(complete_json, 'r', encoding='utf-8') as f:
                            pptx_data = json.load(f)
                        
                        documents = []
                        for page in pptx_data.get('pages', []):
                            page_num = page['page_number']
                            stage3 = page.get('stage3_vlm', {})
                            text_content = stage3.get('text_combined', '')
                            
                            # Create document for this slide
                            doc = Document(
                                page_content=text_content,
                                metadata={
                                    'source': str(file_path),
                                    'file_type': 'pptx',
                                    'page': page_num,
                                    'extraction_method': 'pptx_ocr_pipeline',
                                    'ocr_engine': pptx_data.get('ocr_engine', 'unknown'),
                                    'has_title': page.get('statistics', {}).get('has_title', False),
                                    'total_images': page.get('statistics', {}).get('total_images', 0),
                                    'avg_ocr_confidence': page.get('statistics', {}).get('avg_ocr_confidence', 0.0)
                                }
                            )
                            documents.append(doc)
                        
                        logger.info("pptx_loaded_from_json", num_pages=len(documents))
                        return documents
                    else:
                        logger.warning("pptx_json_not_found", expected_path=str(complete_json))
                
                # Fallback: Try to load with python-pptx directly (basic text extraction)
                try:
                    from pptx import Presentation
                    prs = Presentation(str(file_path))
                    documents = []
                    
                    for slide_num, slide in enumerate(prs.slides, 1):
                        text_content = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                text_content.append(shape.text)
                        
                        if text_content:
                            doc = Document(
                                page_content='\n\n'.join(text_content),
                                metadata={
                                    'source': str(file_path),
                                    'file_type': 'pptx',
                                    'page': slide_num,
                                    'extraction_method': 'python-pptx'
                                }
                            )
                            documents.append(doc)
                    
                    logger.info("pptx_loaded_with_python_pptx", num_slides=len(documents))
                    return documents
                except ImportError:
                    logger.error("python_pptx_not_available")
                    raise ValueError("PPTX file requires python-pptx or pre-processed JSON")
            
            elif file_ext in ['.png', '.jpg', '.jpeg']:
                # Check if we have pre-processed JSON from OCR pipeline
                if processed_json_dir:
                    json_dir = Path(processed_json_dir)
                    complete_json = json_dir / 'complete_document.json'
                    
                    if complete_json.exists():
                        logger.info("loading_image_from_preprocessed_json", json_file=str(complete_json))
                        with open(complete_json, 'r', encoding='utf-8') as f:
                            data = json.load(f)
                        
                        # Build document from OCR-extracted content
                        pages = data.get('pages', [])
                        if pages:
                            page_data = pages[0]  # Single image = single page
                            text_content = page_data.get('text', '')
                            
                            # Use text_blocks if available, otherwise use text
                            if not text_content and page_data.get('text_blocks'):
                                text_blocks = page_data['text_blocks']
                                text_content = '\n'.join(block.get('text', '') for block in text_blocks if block.get('text'))
                            
                            doc = Document(
                                page_content=text_content,
                                metadata={
                                    'source': str(file_path),
                                    'file_type': 'image',
                                    'extraction_method': page_data.get('extraction_method', 'ocr'),
                                    'ocr_engine': page_data.get('ocr_engine', 'unknown'),
                                    'avg_ocr_confidence': page_data.get('avg_ocr_confidence', 0.0)
                                }
                            )
                            
                            logger.info("image_loaded_from_ocr", 
                                       text_length=len(text_content),
                                       avg_confidence=page_data.get('avg_ocr_confidence', 0))
                            return [doc]
                
                # Fallback: Use vision model for images (if no pre-processed JSON)
                return self._load_image_document(file_path)
            else:
                raise ValueError(f"Unsupported file format: {file_ext}")
            
            # For non-PDF files, use the standard loading
            documents = loader.load()
            
            # Log detailed information about loaded documents
            logger.info(
                "document_loaded",
                file_path=str(file_path),
                file_ext=file_ext,
                num_documents=len(documents)
            )
            
            # Log content of each document
            for idx, doc in enumerate(documents):
                content_length = len(doc.page_content) if doc.page_content else 0
                logger.info(
                    "document_page_info",
                    page_index=idx,
                    content_length=content_length,
                    has_content=bool(doc.page_content and doc.page_content.strip()),
                    content_preview=doc.page_content[:200] if doc.page_content else "EMPTY"
                )
            
            return documents
        
        except Exception as e:
            logger.error("document_load_failed", error=str(e), file_path=str(file_path))
            raise
    
    def _load_from_complete_json(self, json_path: Path) -> List[Document]:
        """
        Load document from pre-processed complete_document.json
        This avoids redundant VLM calls since OCR pipeline already extracted everything
        
        Args:
            json_path: Path to complete_document.json
        
        Returns:
            List of LangChain Document objects
        """
        try:
            logger.info("📖 Reading pre-processed JSON (skipping VLM)", json_path=str(json_path))
            
            with open(json_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Handle new format {"pages": [...]} vs old format [...]
            if isinstance(data, dict) and 'pages' in data:
                pages_data = data['pages']
            elif isinstance(data, list):
                pages_data = data
            else:
                raise ValueError(f"Unknown JSON structure in {json_path}")
            
            logger.info("✅ Successfully loaded complete_document.json", num_pages=len(pages_data))
            
            # Load OCR metadata from complete_adaptive_ocr.json for bbox data (for MCP tool)
            ocr_metadata_path = json_path.parent / "complete_adaptive_ocr.json"
            ocr_pages_map = {}
            if ocr_metadata_path.exists():
                try:
                    with open(ocr_metadata_path, 'r', encoding='utf-8') as f:
                        ocr_meta = json.load(f)
                    for page in ocr_meta.get('pages', []):
                        page_num = page.get('page_number')
                        ocr_json_file = page.get('stage1_global', {}).get('ocr_json')
                        if page_num and ocr_json_file:
                            ocr_pages_map[page_num] = ocr_json_file
                    logger.info("✅ Loaded OCR metadata for MCP tool", num_pages=len(ocr_pages_map))
                except Exception as e:
                    logger.warning("failed_to_load_ocr_metadata", error=str(e))
            
            documents = []
            for page_data in pages_data:
                page_num = page_data.get('page_number', len(documents) + 1)
                
                # Handle new format content extraction
                if 'content' in page_data and 'full_text_cleaned' in page_data['content']:
                    page_content = page_data['content']['full_text_cleaned']
                else:
                    # Fallback to old flattening logic
                    page_content = self._flatten_to_searchable_text(page_data)
                
                # Load OCR data for MCP tool (will be serialized to JSON string in vector_store)
                ocr_data = {}
                if page_num in ocr_pages_map:
                    ocr_json_path = json_path.parent / ocr_pages_map[page_num]
                    if ocr_json_path.exists():
                        try:
                            with open(ocr_json_path, 'r', encoding='utf-8') as f:
                                ocr_json = json.load(f)
                            ocr_data = {
                                'text_blocks': ocr_json.get('text_blocks', []),
                                'image_size': ocr_json.get('image_size', {}),
                                'file': ocr_json.get('file', ''),
                                'status': ocr_json.get('status', '')
                            }
                            logger.debug(f"✅ Loaded OCR data for page {page_num}", 
                                       num_blocks=len(ocr_data.get('text_blocks', [])))
                        except Exception as e:
                            logger.warning(f"failed_to_load_ocr_json_for_page_{page_num}", error=str(e))
                
                # Merge metadata
                metadata = {
                        'page': page_num,
                        'page_number': page_num,
                        'page_json': page_data,  # Will be serialized to page_json_raw
                        'ocr_data': ocr_data,    # Will be serialized to ocr_data_raw
                    'page_type': page_data.get('metadata', {}).get('page_type', 'mixed'),
                    'extraction_method': page_data.get('metadata', {}).get('extraction_method', 'vlm_refined')
                    }
                
                # Add VLM refined flag
                if page_data.get('metadata', {}).get('vlm_refined'):
                    metadata['vlm_refined'] = True
                
                # Create Document
                doc = Document(
                    page_content=page_content,
                    metadata=metadata
                )
                
                documents.append(doc)
                
                logger.info(
                    f"✅ Page {page_num} converted to searchable text",
                    page_number=page_num,
                    content_length=len(page_content),
                    has_content=bool(page_content.strip()),
                    has_ocr_data=bool(ocr_data)
                )
            
            return documents
        
        except Exception as e:
            logger.error("failed_to_load_complete_json", error=str(e), json_path=str(json_path))
            raise
    
    def _load_image_document(self, image_path: Path) -> List[Document]:
        """
        Load image document using vision model
        
        Args:
            image_path: Path to image file
        
        Returns:
            List containing single Document with extracted text
        """
        if self.vision_model is None or not self.vision_model.enabled:
            raise ValueError("Vision model is not enabled. Cannot process image files.")
        
        result = self.vision_model.extract_text_from_image(str(image_path))
        
        if result.get('error'):
            raise ValueError(f"Image extraction failed: {result['error']}")
        
        doc = Document(
            page_content=result['text'],
            metadata={
                'source': str(image_path),
                'file_type': 'image',
                'extraction_method': 'vision_model'
            }
        )
        
        return [doc]
    
    def _process_pdf_with_vision(self, pdf_path: Path) -> List[Document]:
        """
        Convert PDF pages to images and extract text using vision model
        
        Args:
            pdf_path: Path to PDF file
        
        Returns:
            List of Documents with extracted text from each page
        """
        if not self.vision_model or not self.vision_model.enabled:
            logger.warning("vision_model_not_enabled")
            return []
        
        logger.info("converting_pdf_to_images", pdf_path=str(pdf_path))
        
        # Convert PDF to images
        images = convert_from_path(str(pdf_path), dpi=200)
        logger.info("pdf_converted_to_images", num_pages=len(images))
        
        documents = []
        temp_files = []
        
        try:
            # Process each page
            for page_num, image in enumerate(images):
                logger.info("processing_pdf_page_with_vision", page=page_num + 1, total=len(images))
                
                # Save image to temp file
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                    image.save(tmp_file.name, 'PNG')
                    temp_files.append(tmp_file.name)
                    
                    # Extract text using vision model
                    result = self.vision_model.extract_text_from_image(tmp_file.name)
                    
                    if result.get('text'):
                        doc = Document(
                            page_content=result['text'],
                            metadata={
                                'source': str(pdf_path),
                                'page': page_num + 1,
                                'file_type': 'pdf',
                                'extraction_method': 'vision_model'
                            }
                        )
                        documents.append(doc)
                        
                        logger.info(
                            "pdf_page_extracted",
                            page=page_num + 1,
                            text_length=len(result['text'])
                        )
                    else:
                        logger.warning(
                            "pdf_page_no_text",
                            page=page_num + 1,
                            error=result.get('error', 'No text extracted')
                        )
        
        finally:
            # Clean up temp files
            for temp_file in temp_files:
                try:
                    os.unlink(temp_file)
                except Exception as e:
                    logger.warning("temp_file_cleanup_failed", file=temp_file, error=str(e))
        
        return documents
    
    def extract_metadata(self, file_path: str) -> Dict[str, Any]:
        """
        Extract metadata from file
        
        Args:
            file_path: Path to file
        
        Returns:
            Dictionary containing file metadata
        """
        file_path = Path(file_path)
        
        # Basic file info
        stat = file_path.stat()
        
        # Calculate checksum
        with open(file_path, 'rb') as f:
            file_hash = hashlib.sha256(f.read()).hexdigest()
        
        metadata = {
            'filename': file_path.name,
            'filepath': str(file_path.absolute()),
            'file_type': file_path.suffix.lower().lstrip('.'),
            'file_size': stat.st_size,
            'created_at': datetime.fromtimestamp(stat.st_ctime).isoformat(),
            'updated_at': datetime.fromtimestamp(stat.st_mtime).isoformat(),
            'checksum': file_hash,
        }
        
        return metadata
    
    def detect_page_content_type(self, page_content: str) -> str:
        """
        Detect page content type
        
        Args:
            page_content: Page text content
        
        Returns:
            Content type (text/table/drawing/mixed)
        """
        if not page_content or not page_content.strip():
            return 'text'
            
        page_content_lower = page_content.lower()
        
        # Check for drawing indicators
        drawing_indicators = [
            'drawing', 'dwg', 'sheet', '图纸', '图号', 'scale', 'rev.', 
            'p&id', 'flow diagram', 'layout', 'datasheet', 'specification'
        ]
        has_drawing = any(indicator in page_content_lower for indicator in drawing_indicators)
        
        # Check for table indicators
        # 1. Explicit keywords
        table_keywords = ['table', '表格', 'list', 'schedule', 'parameters', 'properties', 'characteristics']
        has_table_keyword = any(keyword in page_content_lower for keyword in table_keywords)
        
        # 2. Structure indicators (lines with separators or aligned columns)
        lines = page_content.split('\n')
        table_lines = 0
        for line in lines[:20]:  # Check first 20 lines
            if '|' in line or '\t' in line or '  ' in line:  # Simple heuristic
                table_lines += 1
        
        has_table_structure = table_lines > 3
        
        # Check content length (short content might be scanned or just a title block)
        is_short = len(page_content.strip()) < 200
        
        if has_drawing:
            return 'drawing'
        elif has_table_keyword or has_table_structure:
            return 'table'
        elif is_short:
             # If short but has drawing keywords, it's drawing. Otherwise maybe just text.
             # But here we prioritize "drawing" if it looks like a technical doc page
             return 'mixed'
        else:
            return 'text'
    
    def _flatten_page_json(self, page_json: Dict[str, Any]) -> Dict[str, Any]:
        """
        Flatten page JSON for searchable fields
        Supports both refine_with_vlm.py and vlm_extractor.py formats
        
        Args:
            page_json: Page data from VLM extraction
        
        Returns:
            Flattened data dictionary
        """
        flattened = {}
        
        # Extract visual description from page_analysis (for visual content search)
        page_analysis = page_json.get('page_analysis', {})
        if isinstance(page_analysis, dict):
            # New format: visual_description (merged field)
            visual_desc = page_analysis.get('visual_description', '')
            
            # Legacy format compatibility: merge page_description + layout_structure
            if not visual_desc:
                page_desc = page_analysis.get('page_description', '')
                layout_struct = page_analysis.get('layout_structure', '')
                visual_desc = f"{page_desc} {layout_struct}".strip()
            
            flattened['visual_description'] = visual_desc
            flattened['page_type'] = page_analysis.get('page_type', '')
        else:
            flattened['visual_description'] = ''
            flattened['page_type'] = ''
        
        # Check format and extract accordingly
        if 'content' in page_json:
            # refine_with_vlm.py format
            content = page_json.get('content', {})
            metadata = page_json.get('metadata', {})
            
            # Basic info from metadata
            if isinstance(metadata, dict):
                flattened['drawing_number'] = metadata.get('document_id', '')
                flattened['project_name'] = metadata.get('title', '')
            else:
                flattened['drawing_number'] = ''
                flattened['project_name'] = ''
            
            # Extract from key_fields (only if content is dict)
            if isinstance(content, dict):
                key_fields = content.get('key_fields', [])
                for field in key_fields:
                    field_name = field.get('field', '')
                    field_value = field.get('value', '')
                    if 'number' in field_name.lower() or 'id' in field_name.lower():
                        if not flattened['drawing_number']:
                            flattened['drawing_number'] = field_value
            
            # Equipment and components (empty for non-industrial docs)
            flattened['equipment_tags'] = []
            flattened['equipment_names'] = []
            flattened['all_components'] = ''
            flattened['component_details'] = []
            
            # Table cells from tables (only if content is dict)
            table_cells = []
            if isinstance(content, dict):
                for table in content.get('tables', []):
                    if isinstance(table, dict):
                        # Add description if available
                        if 'description' in table:
                            table_cells.append(table['description'])
            flattened['table_cells'] = table_cells
            
            # All text tokens
            if isinstance(content, dict):
                full_text = content.get('full_text_cleaned', '') or content.get('full_text_raw', '')
            else:
                # content is already a string (from process_pdf_vlm.py)
                full_text = str(content) if content else ''
            flattened['all_text_tokens'] = full_text
        
        else:
            # Original vlm_extractor.py format
            # Document info
            doc_info = page_json.get('document_info', {})
            flattened['drawing_number'] = doc_info.get('drawing_number', '')
            flattened['project_name'] = doc_info.get('project_name', '')
            
            # Equipment tags (exact match)
            equipment = page_json.get('equipment', []) or []
            flattened['equipment_tags'] = [e.get('tag') or e.get('id') for e in equipment if e.get('tag') or e.get('id')]
            flattened['equipment_names'] = [e.get('name') for e in equipment if e.get('name')]
            
            # All components (for search)
            components = page_json.get('components', []) or []
            component_ids = [c.get('id') for c in components if c.get('id')]
            
            # Also get from all_components_list if available
            if 'all_components_list' in page_json:
                component_ids.extend(page_json['all_components_list'])
            
            # Remove duplicates
            component_ids = list(set(component_ids))
            flattened['all_components'] = ' '.join(component_ids)
            
            # Component details (nested)
            flattened['component_details'] = [
                {
                    'id': c.get('id', ''),
                    'type': c.get('type', ''),
                    'value': c.get('value', ''),
                    'reference': c.get('id', ''),
                    'position': c.get('position', '')
                }
                for c in components
            ]
            
            # Flatten tables
            tables = page_json.get('tables', []) or []
            table_cells = []
            for table in tables:
                # Add headers
                headers = table.get('headers', [])
                table_cells.extend([str(h) for h in headers])
                
                # Add all cells from rows
                for row in table.get('rows', []):
                    table_cells.extend([str(cell) for cell in row])
            
            flattened['table_cells'] = table_cells
            
            # All text tokens
            all_texts = page_json.get('all_text', []) or []
            if isinstance(all_texts, list):
                flattened['all_text_tokens'] = ' '.join(all_texts)
            else:
                flattened['all_text_tokens'] = str(all_texts)
        
        return flattened
    
    def _flatten_to_searchable_text(self, page_json: Dict[str, Any]) -> str:
        """
        Convert page JSON to searchable text
        Supports two JSON formats:
        1. refine_with_vlm.py format (content.full_text_cleaned)
        2. vlm_extractor.py format (document_info, equipment, components)
        
        Args:
            page_json: Page data from VLM extraction
        
        Returns:
            Flattened text for content field
        """
        parts = []
        
        # Check if this is refine_with_vlm.py format (has 'content' key)
        if 'content' in page_json:
            content = page_json['content']
            
            # Use full_text_cleaned as primary content
            if 'full_text_cleaned' in content and content['full_text_cleaned']:
                parts.append(content['full_text_cleaned'])
            elif 'full_text_raw' in content and content['full_text_raw']:
                parts.append(content['full_text_raw'])
            
            # Add key fields
            for field in content.get('key_fields', []):
                field_name = field.get('field', '')
                field_value = field.get('value', '')
                if field_name and field_value:
                    parts.append(f"{field_name}: {field_value}")
            
            # Add page description
            if 'page_analysis' in page_json:
                page_desc = page_json['page_analysis'].get('page_description', '')
                if page_desc:
                    parts.append(f"Page Description: {page_desc}")
        
        else:
            # Original vlm_extractor.py format
            # Document info
            doc_info = page_json.get('document_info', {})
            for key, value in doc_info.items():
                if value:
                    parts.append(f"{key}: {value}")
            
            # Equipment
            for equip in page_json.get('equipment', []):
                tag = equip.get('tag') or equip.get('id')
                name = equip.get('name')
                if tag:
                    parts.append(f"Equipment {tag}: {name or ''}")
            
            # Components
            for comp in page_json.get('components', []):
                comp_id = comp.get('id')
                comp_type = comp.get('type')
                value = comp.get('value')
                if comp_id:
                    parts.append(f"Component {comp_id} ({comp_type}): {value or ''}")
            
            # All text
            all_text = page_json.get('all_text', [])
            if isinstance(all_text, list):
                parts.extend(all_text)
            elif isinstance(all_text, str):
                parts.append(all_text)
            
            # Notes
            notes = page_json.get('notes', [])
            if isinstance(notes, list):
                parts.extend(notes)
        
        return '\n'.join(filter(None, parts))
    
    def _convert_pdf_page_to_image(self, pdf_path: Path, page_num: int) -> str:
        """
        Convert single PDF page to image
        
        Args:
            pdf_path: Path to PDF file
            page_num: Page number (0-based)
        
        Returns:
            Path to temporary image file
        """
        if not PDF2IMAGE_AVAILABLE:
            raise ImportError("pdf2image not available")
        
        # Convert single page
        # DPI 150 is sufficient for text recognition and faster than 200
        images = convert_from_path(str(pdf_path), first_page=page_num + 1, last_page=page_num + 1, dpi=150)
        
        if not images:
            raise ValueError(f"Failed to convert page {page_num}")
        
        # Save to temp file
        temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        images[0].save(temp_file.name, 'PNG')
        
        return temp_file.name
    
    def process_document(
        self,
        file_path: str,
        additional_metadata: Optional[Dict[str, Any]] = None,
        processed_json_dir: Optional[str] = None
    ) -> List[Document]:
        """
        Process document: load, split, and add metadata
        
        Args:
            file_path: Path to document file
            additional_metadata: Additional metadata to add
            processed_json_dir: Optional path to directory with pre-processed JSON
        
        Returns:
            List of processed Document chunks
        """
        # Load document (will use pre-processed JSON if available)
        logger.info("loading_document", file_path=file_path, has_processed_json=bool(processed_json_dir))
        documents = self.load_document(file_path, processed_json_dir=processed_json_dir)
        logger.info("document_loaded", num_raw_documents=len(documents))
        
        # Log raw document content
        for i, doc in enumerate(documents):
            logger.debug(
                "raw_document_content",
                doc_index=i,
                content_length=len(doc.page_content) if doc.page_content else 0,
                content_type=type(doc.page_content).__name__,
                content_preview=doc.page_content[:200] if doc.page_content else "EMPTY"
            )
        
        # Extract metadata
        file_metadata = self.extract_metadata(file_path)
        logger.info("metadata_extracted", metadata_keys=list(file_metadata.keys()))
        
        # Merge additional metadata
        if additional_metadata:
            file_metadata.update(additional_metadata)
            logger.info("additional_metadata_merged", keys=list(additional_metadata.keys()))
        
        # Add page-level metadata and flatten JSON
        document_name = Path(file_path).name
        total_pages = len(documents)
        
        for i, doc in enumerate(documents):
            doc.metadata.update(file_metadata)
            
            # Add page metadata
            doc.metadata['page_number'] = i + 1
            doc.metadata['total_pages'] = total_pages
            doc.metadata['document_name'] = document_name
            
            # Flatten page_json if exists
            if 'page_json' in doc.metadata:
                flattened = self._flatten_page_json(doc.metadata['page_json'])
                doc.metadata.update(flattened)
                logger.debug(
                    "page_json_flattened",
                    page_number=i + 1,
                    num_components=len(flattened.get('component_details', [])),
                    num_equipment=len(flattened.get('equipment_tags', []))
                )
        
        # Split documents into chunks (don't split across pages)
        logger.info("splitting_documents", num_documents=len(documents))
        chunks = []
        
        # Determine if we should use page-level chunking
        page_level_indexing = self.config.get('page_level_indexing', True)
        max_page_size = self.config.get('max_page_size_chars', 4000)
        
        for doc in documents:
            # Check document length
            content_len = len(doc.page_content) if doc.page_content else 0
            
            if page_level_indexing and content_len < max_page_size:
                # Keep as single chunk
                doc.metadata['chunk_type'] = 'full_page'
                # Ensure content is string
                if not isinstance(doc.page_content, str):
                    doc.page_content = str(doc.page_content) if doc.page_content else ""
                chunks.append(doc)
            else:
                # Split large pages
                page_chunks = self.text_splitter.split_documents([doc])
                for i, chunk in enumerate(page_chunks):
                    chunk.metadata['chunk_type'] = 'part_page'
                    chunk.metadata['part_index'] = i
                chunks.extend(page_chunks)
                
        logger.info("documents_split", num_chunks=len(chunks))
        
        # Filter out empty chunks and validate content
        valid_chunks = []
        skipped_count = 0
        for idx, chunk in enumerate(chunks):
            # Ensure page_content is string and not empty
            if not isinstance(chunk.page_content, str):
                logger.warning(
                    "converting_non_string_content",
                    chunk_index=idx,
                    original_type=type(chunk.page_content).__name__
                )
                chunk.page_content = str(chunk.page_content)
            
            # Skip empty or whitespace-only chunks
            if chunk.page_content.strip():
                valid_chunks.append(chunk)
                logger.debug(
                    "valid_chunk",
                    chunk_index=idx,
                    content_length=len(chunk.page_content),
                    content_preview=chunk.page_content[:100]
                )
            else:
                skipped_count += 1
                logger.warning("skipping_empty_chunk", chunk_index=idx)
        
        if skipped_count > 0:
            logger.warning("empty_chunks_skipped", count=skipped_count)
        
        # Add chunk metadata
        for i, chunk in enumerate(valid_chunks):
            chunk.metadata['chunk_index'] = i
            chunk.metadata['total_chunks'] = len(valid_chunks)
            chunk.metadata['document_id'] = file_metadata['checksum']
            chunk.metadata['chunk_id'] = f"{file_metadata['checksum']}_chunk_{i}"
        
        logger.info(
            "document_processed",
            file_path=file_path,
            num_chunks=len(valid_chunks),
            total_size=file_metadata['file_size']
        )
        
        return valid_chunks
    
    def process_zip(
        self,
        zip_path: str,
        extract_dir: Optional[str] = None,
        additional_metadata: Optional[Dict[str, Any]] = None
    ) -> List[Document]:
        """
        Process ZIP file containing multiple documents
        
        Args:
            zip_path: Path to ZIP file
            extract_dir: Directory to extract files (temp dir if None)
            additional_metadata: Additional metadata to add
        
        Returns:
            List of all processed document chunks
        """
        if extract_dir is None:
            extract_dir = f"./uploads/extracted_{Path(zip_path).stem}"
        
        extract_path = Path(extract_dir)
        extract_path.mkdir(parents=True, exist_ok=True)
        
        all_chunks = []
        
        try:
            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                zip_ref.extractall(extract_path)
                
                logger.info("zip_extracted", zip_path=zip_path, extract_dir=str(extract_path))
                
                # Process each file in ZIP
                for root, _, files in os.walk(extract_path):
                    for file in files:
                        file_path = Path(root) / file
                        
                        # Skip hidden files and unsupported formats
                        if file.startswith('.'):
                            continue
                        
                        if file_path.suffix.lower().lstrip('.') not in self.config.get('supported_formats', []):
                            logger.debug("file_skipped", file_path=str(file_path))
                            continue
                        
                        try:
                            chunks = self.process_document(
                                str(file_path),
                                additional_metadata
                            )
                            all_chunks.extend(chunks)
                        except Exception as e:
                            logger.error(
                                "file_processing_failed",
                                error=str(e),
                                file_path=str(file_path)
                            )
                            continue
        
        except Exception as e:
            logger.error("zip_processing_failed", error=str(e), zip_path=zip_path)
            raise
        
        logger.info(
            "zip_processed",
            zip_path=zip_path,
            total_chunks=len(all_chunks),
            num_files=len(set(c.metadata['filepath'] for c in all_chunks))
        )
        
        return all_chunks
    
    def process_batch(
        self,
        file_paths: List[str],
        additional_metadata: Optional[Dict[str, Any]] = None
    ) -> List[Document]:
        """
        Process multiple documents in batch
        
        Args:
            file_paths: List of file paths
            additional_metadata: Additional metadata to add
        
        Returns:
            List of all processed document chunks
        """
        all_chunks = []
        
        for file_path in file_paths:
            try:
                chunks = self.process_document(file_path, additional_metadata)
                all_chunks.extend(chunks)
            except Exception as e:
                logger.error(
                    "batch_file_processing_failed",
                    error=str(e),
                    file_path=file_path
                )
                continue
        
        logger.info("batch_processed", total_files=len(file_paths), total_chunks=len(all_chunks))
        
        return all_chunks
    
    def _split_word_by_pages(self, file_path: Path) -> List[Document]:
        """
        Split Word document by page breaks
        
        Args:
            file_path: Path to Word document
        
        Returns:
            List of Documents, one per page
        """
        try:
            from docx import Document as DocxDocument
        except ImportError:
            logger.warning("python-docx_not_available")
            raise
        
        doc = DocxDocument(str(file_path))
        pages = []
        current_page_text = []
        
        for para in doc.paragraphs:
            # Check for page break
            if '\f' in para.text or '\x0c' in para.text:
                # Save current page
                if current_page_text:
                    page_content = '\n'.join(current_page_text)
                    pages.append(page_content)
                    current_page_text = []
            else:
                if para.text.strip():
                    current_page_text.append(para.text)
        
        # Add last page
        if current_page_text:
            pages.append('\n'.join(current_page_text))
        
        # If no page breaks found, treat as single page
        if not pages:
            pages = ['\n'.join(para.text for para in doc.paragraphs if para.text.strip())]
        
        # Create Document objects
        documents = []
        for i, page_text in enumerate(pages):
            doc = Document(
                page_content=page_text,
                metadata={
                    'source': str(file_path),
                    'page': i + 1,
                    'file_type': 'docx',
                    'extraction_method': 'text',
                    'page_type': 'text'
                }
            )
            documents.append(doc)
        
        logger.info("word_split_by_pages", num_pages=len(documents))
        
        return documents
    
    def _process_excel_sheets(self, file_path: Path) -> List[Document]:
        """
        Process Excel file by sheets (each sheet as one page)
        
        Args:
            file_path: Path to Excel file
        
        Returns:
            List of Documents, one per sheet
        """
        try:
            import pandas as pd
        except ImportError:
            logger.warning("pandas_not_available")
            raise
        
        # Read all sheets
        excel_file = pd.ExcelFile(str(file_path))
        documents = []
        
        for i, sheet_name in enumerate(excel_file.sheet_names):
            df = excel_file.parse(sheet_name)
            
            # Convert dataframe to text
            sheet_text = f"Sheet: {sheet_name}\n\n"
            sheet_text += df.to_string(index=False)
            
            doc = Document(
                page_content=sheet_text,
                metadata={
                    'source': str(file_path),
                    'page': i + 1,
                    'sheet_name': sheet_name,
                    'file_type': 'xlsx',
                    'extraction_method': 'text',
                    'page_type': 'table'
                }
            )
            documents.append(doc)
        
        logger.info("excel_processed_by_sheets", num_sheets=len(documents))
        
        return documents

